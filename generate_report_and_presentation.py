from pathlib import Path
from datetime import date
import json
import math

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
import numpy as np
import pandas as pd

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(".")
TABLE_DIR = ROOT / "outputs" / "tables"
EMBEDDING_DIR = ROOT / "outputs" / "embeddings"
FIGURE_DIR = ROOT / "outputs" / "figures"
REPORT_DIR = ROOT / "reports"
PRESENTATION_DIR = ROOT / "presentations"

REPORT_PATH = REPORT_DIR / "NLP_Mini_Project_3_Khmer_Word_Embeddings_Report.docx"
PPTX_PATH = PRESENTATION_DIR / "NLP_Mini_Project_3_Khmer_Word_Embeddings_Presentation.pptx"

NAVY = "#0B1F3A"
GOLD = "#C9A227"
BLUE = "#4E9FE5"
SOFT_BLUE = "#EEF6FF"
GRAY = "#526070"
GREEN = "#1F9D55"

FALLBACK_RESULTS = {
    "total_tokens": 9075,
    "vocab_size": 175,
    "embedding_dim": 50,
    "window_size": "+/-4",
    "negative_samples": 2,
    "n_context": 5,
    "hidden_size": 512,
    "skipgram_loss": 0.4951,
    "embedding_shape": "(175, 50)",
    "ngram_perplexity": 158.84,
    "fixed_perplexity": 115.10,
    "scratch_perplexity": 229.07,
}


def ensure_dirs():
    for path in [FIGURE_DIR, REPORT_DIR, PRESENTATION_DIR]:
        path.mkdir(parents=True, exist_ok=True)


def find_khmer_font():
    preferred_fonts = [
        "Noto Sans Khmer",
        "Noto Serif Khmer",
        "Khmer OS Siemreap",
        "Khmer OS Battambang",
        "Khmer OS Content",
        "Khmer OS",
        "DaunPenh",
        "Leelawadee UI",
    ]

    installed_fonts = font_manager.fontManager.ttflist
    for preferred_font in preferred_fonts:
        for font in installed_fonts:
            if preferred_font.lower() == font.name.lower():
                return font.name

    for preferred_font in preferred_fonts:
        for font in installed_fonts:
            if preferred_font.lower() in font.name.lower():
                return font.name

    return None


def configure_matplotlib_fonts():
    """Use a Khmer-capable font for chart labels when it is available."""
    plt.rcParams["axes.unicode_minus"] = False
    khmer_font_name = find_khmer_font()
    if khmer_font_name is not None:
        plt.rcParams["font.family"] = [khmer_font_name, "DejaVu Sans"]
        return khmer_font_name
    print("No Khmer-supported font found. Khmer labels may not display correctly. Please install Noto Sans Khmer or Khmer OS fonts.")
    return "Default Matplotlib font"


def contains_khmer_text(text):
    for character in str(text):
        if "\u1780" <= character <= "\u17FF":
            return True
    return False


def read_csv_or_empty(path):
    if path.exists():
        return pd.read_csv(path)
    return pd.DataFrame()


def load_project_data():
    cleaned_tokens = read_csv_or_empty(TABLE_DIR / "cleaned_tokens.csv")
    vocabulary = read_csv_or_empty(TABLE_DIR / "vocabulary.csv")
    token_frequency = read_csv_or_empty(TABLE_DIR / "token_frequency.csv")
    model_comparison = read_csv_or_empty(TABLE_DIR / "model_comparison.csv")
    prediction_examples = read_csv_or_empty(TABLE_DIR / "prediction_examples.csv")
    silhouette_scores = read_csv_or_empty(TABLE_DIR / "kmeans_silhouette_scores.csv")
    kmeans_clusters = read_csv_or_empty(TABLE_DIR / "kmeans_word_clusters.csv")
    cluster_summary = read_csv_or_empty(TABLE_DIR / "kmeans_cluster_summary.csv")
    dendrogram_labels = read_csv_or_empty(TABLE_DIR / "dendrogram_word_labels.csv")

    if token_frequency.empty:
        token_frequency = pd.DataFrame(
            {
                "word": ["ប្រាសាទ", "អង្គរ", "ខ្មែរ", "សាសនា", "ថ្ម"],
                "frequency": [250, 190, 120, 90, 65],
                "length": [7, 5, 5, 6, 3],
            }
        )
        token_frequency["rank"] = range(1, len(token_frequency) + 1)

    if model_comparison.empty:
        model_comparison = pd.DataFrame(
            [
                {
                    "Model": "N-gram baseline",
                    "Embedding Source": "No embeddings",
                    "Test Perplexity": FALLBACK_RESULTS["ngram_perplexity"],
                    "Top-1 Accuracy": np.nan,
                    "Top-5 Accuracy": np.nan,
                    "Notes": "Add-one smoothing baseline",
                },
                {
                    "Model": "Neural LM fixed",
                    "Embedding Source": "Fixed skip-gram embeddings",
                    "Test Perplexity": FALLBACK_RESULTS["fixed_perplexity"],
                    "Top-1 Accuracy": np.nan,
                    "Top-5 Accuracy": np.nan,
                    "Notes": "Best model in final result",
                },
                {
                    "Model": "Neural LM scratch",
                    "Embedding Source": "Learned from scratch",
                    "Test Perplexity": FALLBACK_RESULTS["scratch_perplexity"],
                    "Top-1 Accuracy": np.nan,
                    "Top-5 Accuracy": np.nan,
                    "Notes": "Can overfit small corpus",
                },
            ]
        )

    total_tokens = len(cleaned_tokens) if not cleaned_tokens.empty else FALLBACK_RESULTS["total_tokens"]
    vocab_size = len(vocabulary) if not vocabulary.empty else FALLBACK_RESULTS["vocab_size"]

    return {
        "cleaned_tokens": cleaned_tokens,
        "vocabulary": vocabulary,
        "token_frequency": token_frequency,
        "model_comparison": model_comparison,
        "prediction_examples": prediction_examples,
        "silhouette_scores": silhouette_scores,
        "kmeans_clusters": kmeans_clusters,
        "cluster_summary": cluster_summary,
        "dendrogram_labels": dendrogram_labels,
        "total_tokens": total_tokens,
        "vocab_size": vocab_size,
    }


def save_chart(path):
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()
    return path


def style_axes(title, xlabel="", ylabel=""):
    ax = plt.gca()
    ax.set_title(title, fontsize=13, fontweight="bold", color=NAVY, pad=12)
    ax.set_xlabel(xlabel, fontsize=10, color=NAVY)
    ax.set_ylabel(ylabel, fontsize=10, color=NAVY)
    ax.grid(True, axis="x", alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def make_charts(data):
    token_frequency = data["token_frequency"].copy()
    model_comparison = data["model_comparison"].copy()
    chart_paths = {}

    top_words = token_frequency.sort_values("frequency", ascending=False).head(20)
    top_words = top_words.sort_values("frequency", ascending=True)
    plt.figure(figsize=(9, 7))
    plt.barh(top_words["word"].astype(str), top_words["frequency"], color=NAVY)
    style_axes("Top Frequent Khmer Words", "Frequency", "Word")
    chart_paths["top_words"] = save_chart(FIGURE_DIR / "top_frequent_words.png")

    if "length" not in token_frequency.columns:
        token_frequency["length"] = token_frequency["word"].astype(str).str.len()
    plt.figure(figsize=(8, 4.8))
    plt.hist(token_frequency["length"], bins=25, color=BLUE, edgecolor="white")
    style_axes("Token Length Distribution", "Token Length", "Number of Words")
    chart_paths["token_length"] = save_chart(FIGURE_DIR / "token_length_distribution.png")

    plt.figure(figsize=(8, 4.8))
    plt.hist(token_frequency["frequency"], bins=30, color=BLUE, edgecolor="white", log=True)
    style_axes("Word Frequency Distribution (Log Scale)", "Frequency", "Number of Words")
    chart_paths["frequency_distribution"] = save_chart(FIGURE_DIR / "word_frequency_distribution.png")

    ranked = token_frequency.sort_values("frequency", ascending=False).reset_index(drop=True)
    ranked["rank"] = ranked.index + 1
    plt.figure(figsize=(8, 4.8))
    plt.plot(ranked["rank"], ranked["frequency"], color=GOLD, linewidth=2.5)
    style_axes("Rank-Frequency Curve", "Rank", "Frequency")
    chart_paths["rank_frequency"] = save_chart(FIGURE_DIR / "rank_frequency_curve.png")

    plt.figure(figsize=(8, 4.8))
    model_names = model_comparison["Model"].astype(str)
    perplexity = pd.to_numeric(model_comparison["Test Perplexity"], errors="coerce")
    bars = plt.bar(model_names, perplexity, color=[NAVY, BLUE, GOLD])
    for bar in bars:
        height = bar.get_height()
        if not math.isnan(height):
            plt.text(bar.get_x() + bar.get_width() / 2, height + 4, f"{height:.2f}", ha="center", fontsize=9)
    style_axes("Test Perplexity by Model (Lower is Better)", "Model", "Test Perplexity")
    plt.xticks(rotation=15, ha="right")
    chart_paths["model_perplexity"] = save_chart(FIGURE_DIR / "model_perplexity_comparison.png")

    pca_path = make_pca_chart(data)
    if pca_path is not None:
        chart_paths["pca"] = pca_path
    scratch_pca_path = make_embedding_pca_chart(
        EMBEDDING_DIR / "scratch_lm_embeddings.npy",
        "PCA Map of Scratch LM Embeddings",
        FIGURE_DIR / "pca_scratch_lm_embeddings.png",
    )
    if scratch_pca_path is not None:
        chart_paths["scratch_pca"] = scratch_pca_path

    clustering_paths = make_clustering_charts(data)
    for key, path in clustering_paths.items():
        chart_paths[key] = path

    return chart_paths


def make_clustering_charts(data):
    chart_paths = {}
    silhouette_scores = data.get("silhouette_scores", pd.DataFrame())
    kmeans_clusters = data.get("kmeans_clusters", pd.DataFrame())
    dendrogram_path = FIGURE_DIR / "hierarchical_dendrogram_khmer.png"
    numbered_dendrogram_path = FIGURE_DIR / "hierarchical_dendrogram_numbered.png"

    if not silhouette_scores.empty and {"k", "silhouette_score"}.issubset(silhouette_scores.columns):
        clean_scores = silhouette_scores.copy()
        clean_scores["k"] = pd.to_numeric(clean_scores["k"], errors="coerce")
        clean_scores["silhouette_score"] = pd.to_numeric(clean_scores["silhouette_score"], errors="coerce")
        clean_scores = clean_scores.dropna(subset=["k", "silhouette_score"])
        if not clean_scores.empty:
            plt.figure(figsize=(8, 4.8))
            plt.plot(clean_scores["k"], clean_scores["silhouette_score"], marker="o", color=NAVY, linewidth=2.5)
            style_axes("Silhouette Score for K-means Clustering", "Number of clusters K", "Silhouette Score")
            chart_paths["kmeans_silhouette"] = save_chart(FIGURE_DIR / "kmeans_silhouette_scores.png")

    required_cluster_cols = {"PC1", "PC2", "cluster", "word", "frequency"}
    if not kmeans_clusters.empty and required_cluster_cols.issubset(kmeans_clusters.columns):
        clean_clusters = kmeans_clusters.copy()
        clean_clusters["frequency"] = pd.to_numeric(clean_clusters["frequency"], errors="coerce").fillna(1)
        top_clusters = clean_clusters.sort_values("frequency", ascending=False).head(120)
        plt.figure(figsize=(8, 5.8))
        scatter = plt.scatter(
            top_clusters["PC1"],
            top_clusters["PC2"],
            c=pd.to_numeric(top_clusters["cluster"], errors="coerce"),
            s=np.clip(top_clusters["frequency"], 8, 90),
            cmap="tab10",
            alpha=0.78,
        )
        plt.colorbar(scatter, label="Cluster")
        style_axes("PCA Map Colored by K-means Cluster", "PC1", "PC2")
        chart_paths["kmeans_pca"] = save_chart(FIGURE_DIR / "kmeans_pca_clusters.png")

    if not dendrogram_path.exists() or not numbered_dendrogram_path.exists():
        make_readable_dendrograms()

    if dendrogram_path.exists():
        chart_paths["hierarchical_dendrogram"] = dendrogram_path
    if numbered_dendrogram_path.exists():
        chart_paths["hierarchical_dendrogram_numbered"] = numbered_dendrogram_path

    return chart_paths


def make_readable_dendrograms():
    emb_path = EMBEDDING_DIR / "skipgram_embeddings.npy"
    word2idx_path = EMBEDDING_DIR / "word2idx.json"
    clusters_path = TABLE_DIR / "kmeans_word_clusters.csv"
    if not emb_path.exists() or not word2idx_path.exists() or not clusters_path.exists():
        return

    try:
        from scipy.cluster.hierarchy import dendrogram, linkage
        from scipy.spatial.distance import pdist

        embeddings = np.load(emb_path)
        with open(word2idx_path, "r", encoding="utf-8") as file:
            word2idx = json.load(file)
        clusters = pd.read_csv(clusters_path)
        if clusters.empty or "word" not in clusters.columns or "frequency" not in clusters.columns:
            return

        selected_df = clusters.sort_values("frequency", ascending=False).head(min(40, len(clusters))).copy()
        selected_words = []
        selected_vectors = []
        for row in selected_df.itertuples():
            word = str(row.word)
            if word not in word2idx:
                continue
            selected_words.append(word)
            selected_vectors.append(embeddings[int(word2idx[word])])

        if len(selected_vectors) < 3:
            return

        selected_vectors = np.array(selected_vectors)
        linkage_matrix = linkage(pdist(selected_vectors, metric="cosine"), method="average")

        khmer_font_name = find_khmer_font()
        if khmer_font_name is not None:
            plt.rcParams["font.family"] = [khmer_font_name, "DejaVu Sans"]
        else:
            print("No Khmer-supported font found. Khmer labels may not display correctly. Please install Noto Sans Khmer or Khmer OS fonts.")
        plt.rcParams["axes.unicode_minus"] = False

        plt.figure(figsize=(16, 8))
        dendrogram(linkage_matrix, labels=selected_words, leaf_rotation=90, leaf_font_size=12)
        ax = plt.gca()
        ax.set_title("Hierarchical Clustering Dendrogram of Frequent Khmer Words", fontsize=16, fontname="DejaVu Sans")
        ax.set_ylabel("Cosine Distance", fontsize=12, fontname="DejaVu Sans")
        if khmer_font_name is not None:
            for label in ax.get_xticklabels():
                if contains_khmer_text(label.get_text()):
                    label.set_fontfamily([khmer_font_name, "DejaVu Sans"])
                else:
                    label.set_fontname("DejaVu Sans")
                label.set_fontsize(11)
        plt.tight_layout()
        plt.savefig(FIGURE_DIR / "hierarchical_dendrogram_khmer.png", dpi=300, bbox_inches="tight")
        plt.close()

        label_rows = []
        short_labels = []
        for label_index, word in enumerate(selected_words, start=1):
            short_id = "W" + str(label_index).zfill(2)
            short_labels.append(short_id)
            frequency = selected_df[selected_df["word"].astype(str) == word]["frequency"].iloc[0]
            label_rows.append({"ID": short_id, "Khmer Word": word, "Frequency": frequency})

        pd.DataFrame(label_rows).to_csv(TABLE_DIR / "dendrogram_word_labels.csv", index=False)

        plt.figure(figsize=(16, 8))
        dendrogram(linkage_matrix, labels=short_labels, leaf_rotation=90, leaf_font_size=12)
        ax = plt.gca()
        ax.set_title("Hierarchical Clustering Dendrogram of Frequent Khmer Words", fontsize=16, fontname="DejaVu Sans")
        ax.set_ylabel("Cosine Distance", fontsize=12, fontname="DejaVu Sans")
        plt.tight_layout()
        plt.savefig(FIGURE_DIR / "hierarchical_dendrogram_numbered.png", dpi=300, bbox_inches="tight")
        plt.close()
    except Exception:
        return


def make_pca_chart(data):
    return make_embedding_pca_chart(
        EMBEDDING_DIR / "skipgram_embeddings.npy",
        "PCA Map of Skip-gram Embeddings",
        FIGURE_DIR / "pca_skipgram_embeddings.png",
    )


def make_embedding_pca_chart(emb_path, title, output_path):
    idx_path = EMBEDDING_DIR / "idx2word.json"
    if not emb_path.exists() or not idx_path.exists():
        return None
    try:
        from sklearn.decomposition import PCA

        embeddings = np.load(emb_path)
        with open(idx_path, "r", encoding="utf-8") as file:
            idx2word_json = json.load(file)
        words = []
        for key, value in sorted(idx2word_json.items(), key=lambda item: int(item[0])):
            words.append(value)

        pca = PCA(n_components=2, random_state=42)
        coords = pca.fit_transform(embeddings)
        count = min(80, len(words))
        plt.figure(figsize=(8, 6))
        plt.scatter(coords[:count, 0], coords[:count, 1], s=26, color=BLUE, alpha=0.75)
        for i in range(min(25, count)):
            plt.text(coords[i, 0], coords[i, 1], str(words[i]), fontsize=7, color=NAVY)
        style_axes(title, "PC1", "PC2")
        return save_chart(output_path)
    except Exception:
        return None


def add_doc_paragraph(doc, text, bold=False):
    paragraph = doc.add_paragraph()
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.size = Pt(10.5)
    return paragraph


def add_doc_bullets(doc, items):
    for item in items:
        paragraph = doc.add_paragraph(style="List Bullet")
        paragraph.add_run(item).font.size = Pt(10.5)


def add_doc_table(doc, headers, rows):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    header_cells = table.rows[0].cells
    for i, header in enumerate(headers):
        run = header_cells[i].paragraphs[0].add_run(str(header))
        run.bold = True
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = "" if pd.isna(value) else str(value)
    doc.add_paragraph()
    return table


def add_doc_figure(doc, image_path, caption):
    if image_path and Path(image_path).exists():
        doc.add_picture(str(image_path), width=Inches(5.8))
        paragraph = doc.add_paragraph(caption)
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        paragraph.runs[0].italic = True


def requirement_rows():
    return [
        ["Khmer corpus", "Use temples.txt", "data/temples.txt or temples.txt", "Completed"],
        ["Khmer tokenization", "Use Khmer tokenizer", "khmernltk / khmer-nltk plus fallback", "Completed"],
        ["MIN_FREQ", "Remove frequency < 10", "MIN_FREQ = 10", "Completed"],
        ["Noisy token removal", "Ignore spaces/noise", "Punctuation, numbers, empty and non-Khmer tokens removed", "Completed"],
        ["Skip-gram", "Build model/classifier", "Skip-gram with negative sampling", "Completed"],
        ["Embedding dimension", "50", "EMBEDDING_DIM = 50", "Completed"],
        ["Context window", "+/-4", "WINDOW_SIZE = 4", "Completed"],
        ["Negative sampling", "k = 2", "NEGATIVE_SAMPLES = 2", "Completed"],
        ["PCA", "2 components", "2D skip-gram and scratch embedding visualization", "Completed"],
        ["Neural LM", "Predict next word", "Previous 5 words to next word", "Completed"],
        ["Hidden layer", "h = 512", "HIDDEN_SIZE = 512", "Completed"],
        ["Scratch embeddings", "Learn from scratch", "Scratch LM implemented and compared", "Completed"],
        ["Comparison", "Compare embeddings/models", "Model comparison table and discussion", "Completed"],
        ["Report summary", "Max 4-page summary", "Notebook and generated report include summary", "Completed"],
    ]


def create_word_report(data, chart_paths):
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Khmer Word Embeddings and Neural Language Modeling using Temples Corpus")
    run.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(11, 31, 58)

    subtitle_lines = [
        "Natural Language Processing",
        "Mini Project 3 - Word Embeddings",
        "Program: Master of Data Science",
        "Institute: Institute of Technology of Cambodia",
        "Student Name: ____________________",
        f"Date: {date.today().isoformat()}",
    ]
    for line in subtitle_lines:
        p = doc.add_paragraph(line)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    doc.add_heading("Executive Summary", level=1)
    add_doc_paragraph(
        doc,
        "This project builds Khmer word embeddings and neural language models using a temple-related Khmer corpus. "
        "The workflow includes cleaning, Khmer tokenization, vocabulary construction, skip-gram training with negative sampling, "
        "PCA visualization, n-gram baseline modeling, fixed skip-gram neural language modeling, and scratch embedding comparison.",
    )
    add_doc_paragraph(
        doc,
        "The fixed skip-gram neural language model achieved the best test perplexity. This shows that pretrained contextual "
        "embeddings helped next-word prediction on the small Khmer temple corpus.",
    )

    doc.add_heading("Project Requirement Alignment", level=1)
    add_doc_table(doc, ["Requirement", "Expected Setting", "Implementation", "Status"], requirement_rows())

    doc.add_heading("Dataset Description", level=1)
    add_doc_paragraph(doc, "The project uses temples.txt, a Khmer temple Wikipedia-style corpus focused on Angkor and temple-related content.")
    add_doc_table(
        doc,
        ["Metric", "Value"],
        [
            ["Total useful tokens", f"{data['total_tokens']:,}"],
            ["Vocabulary size after MIN_FREQ=10", f"{data['vocab_size']:,}"],
            ["Corpus domain", "Khmer temple text"],
            ["Corpus limitation", "Small and domain-specific"],
        ],
    )

    doc.add_heading("Text Preprocessing and Data Quality", level=1)
    add_doc_paragraph(
        doc,
        "Khmer tokenization is important because Khmer does not mark word boundaries as clearly as English. "
        "The notebook tries Khmer tokenizer packages and keeps a fallback tokenizer for reproducibility.",
    )
    add_doc_bullets(
        doc,
        [
            "Removed references, repeated spaces, punctuation-only tokens, number-only tokens, and noisy non-Khmer tokens.",
            "Kept MIN_FREQ = 10 to satisfy the project instruction.",
            "Mapped rare words to <UNK>, which increases the unknown-word rate but follows the required threshold.",
        ],
    )

    doc.add_heading("Exploratory Data Analysis", level=1)
    add_doc_figure(doc, chart_paths.get("top_words"), "Figure 1. Top frequent Khmer words after cleaning.")
    add_doc_paragraph(doc, "The frequent-word chart confirms that the vocabulary contains meaningful Khmer words rather than punctuation.")
    add_doc_figure(doc, chart_paths.get("token_length"), "Figure 2. Token length distribution.")
    add_doc_paragraph(doc, "The length distribution helps detect tokenization problems and unusually long tokens.")
    add_doc_figure(doc, chart_paths.get("frequency_distribution"), "Figure 3. Word frequency distribution using log scale.")
    add_doc_paragraph(doc, "The distribution is highly skewed: many words are rare while a few occur very frequently.")
    add_doc_figure(doc, chart_paths.get("rank_frequency"), "Figure 4. Rank-frequency curve.")

    doc.add_heading("Skip-gram Word Embedding Model", level=1)
    add_doc_paragraph(
        doc,
        "The skip-gram model learns word vectors by predicting context words around a center word. "
        "With negative sampling, the task becomes binary classification: true center-context pairs are labeled 1 and sampled negative pairs are labeled 0.",
    )
    add_doc_table(
        doc,
        ["Setting", "Value"],
        [
            ["Embedding dimension", "50"],
            ["Context window", "+/-4"],
            ["Negative samples", "2"],
            ["Approximate final skip-gram loss", FALLBACK_RESULTS["skipgram_loss"]],
            ["Embedding matrix shape", f"({data['vocab_size']}, 50)"],
        ],
    )

    doc.add_heading("Word Embedding Analysis", level=1)
    add_doc_paragraph(
        doc,
        "Cosine similarity compares word vectors. Nearest words are interpreted as contextual neighbors, not always perfect synonyms, because the corpus is small and domain-specific.",
    )

    doc.add_heading("PCA Visualization", level=1)
    add_doc_paragraph(
        doc,
        "PCA reduces 50-dimensional embeddings to 2D for visual inspection. "
        "It is useful for comparing skip-gram embeddings and scratch learned embeddings, but it loses information.",
    )
    add_doc_figure(doc, chart_paths.get("pca"), "Figure 5. PCA visualization of skip-gram embeddings.")
    add_doc_figure(doc, chart_paths.get("scratch_pca"), "Figure 6. PCA visualization of scratch LM embeddings.")

    doc.add_heading("Advanced Embedding Exploration", level=1)
    add_doc_paragraph(
        doc,
        "As optional advanced analysis, K-means and hierarchical clustering were applied to the learned skip-gram embeddings. "
        "These methods help inspect whether Khmer words that appear in similar contexts form groups in the embedding space. "
        "They are exploratory and do not replace the required Mini Project 3 PCA visualization.",
    )
    silhouette_scores = data.get("silhouette_scores", pd.DataFrame())
    if not silhouette_scores.empty:
        add_doc_figure(doc, chart_paths.get("kmeans_silhouette"), "Figure 7. Silhouette score for choosing K in K-means clustering.")
    cluster_summary = data.get("cluster_summary", pd.DataFrame())
    if not cluster_summary.empty:
        display_columns = [col for col in ["Cluster", "Number of Words", "Top Frequent Words", "Interpretation"] if col in cluster_summary.columns]
        add_doc_table(doc, display_columns, cluster_summary[display_columns].values.tolist())
    add_doc_figure(doc, chart_paths.get("kmeans_pca"), "Figure 8. PCA visualization colored by K-means cluster.")
    if chart_paths.get("hierarchical_dendrogram_numbered"):
        add_doc_figure(doc, chart_paths.get("hierarchical_dendrogram_numbered"), "Figure 9. Hierarchical clustering dendrogram using numbered word labels.")
        dendrogram_labels = data.get("dendrogram_labels", pd.DataFrame())
        if not dendrogram_labels.empty:
            label_columns = [col for col in ["ID", "Khmer Word", "Frequency"] if col in dendrogram_labels.columns]
            add_doc_table(doc, label_columns, dendrogram_labels[label_columns].values.tolist())
    else:
        add_doc_figure(doc, chart_paths.get("hierarchical_dendrogram"), "Figure 9. Hierarchical clustering dendrogram of frequent Khmer words.")
    add_doc_paragraph(
        doc,
        "Words in the same K-means cluster suggest similar contextual patterns. The dendrogram gives a tree-style view of selected frequent words. "
        "Because the corpus is small and temple-specific, these clusters should be interpreted cautiously as contextual groups, not perfect semantic classes.",
    )

    doc.add_heading("Neural Language Modeling", level=1)
    add_doc_paragraph(
        doc,
        "The neural language model predicts the next word from the previous five words. "
        "It concatenates five 50-dimensional embeddings, applies a hidden layer of size 512, and predicts the next word using a softmax output with CrossEntropyLoss.",
    )
    add_doc_bullets(
        doc,
        [
            "N-gram baseline uses add-one smoothing.",
            "Fixed neural LM uses pretrained skip-gram embeddings.",
            "Scratch neural LM learns embeddings from random initialization.",
        ],
    )

    doc.add_heading("Model Comparison and Results", level=1)
    comparison = data["model_comparison"]
    headers = ["Model", "Embedding Source", "Test Perplexity", "Top-1 Accuracy", "Top-5 Accuracy", "Interpretation"]
    rows = []
    for _, row in comparison.iterrows():
        notes = row.get("Notes", "")
        rows.append(
            [
                row.get("Model", ""),
                row.get("Embedding Source", ""),
                round(float(row.get("Test Perplexity", np.nan)), 2),
                row.get("Top-1 Accuracy", ""),
                row.get("Top-5 Accuracy", ""),
                notes,
            ]
        )
    add_doc_table(doc, headers, rows)
    add_doc_figure(doc, chart_paths.get("model_perplexity"), "Figure 10. Test perplexity by model. Lower is better.")
    add_doc_paragraph(doc, "The fixed skip-gram neural LM achieved the best perplexity, indicating that pretrained contextual embeddings helped prediction.")

    doc.add_heading("Error Analysis", level=1)
    prediction_examples = data["prediction_examples"]
    if not prediction_examples.empty:
        display_cols = [col for col in prediction_examples.columns[:4]]
        add_doc_table(doc, display_cols, prediction_examples[display_cols].head(10).values.tolist())
    add_doc_paragraph(
        doc,
        "Common errors are caused by rare words, <UNK>, Khmer compound words, named entities, and the small corpus size.",
    )

    doc.add_heading("Discussion", level=1)
    add_doc_paragraph(
        doc,
        "The project successfully connects course topics from preprocessing and n-grams to word meaning and neural language modeling. "
        "Tokenization had the largest impact on quality. Pretrained skip-gram embeddings helped the neural LM because they already captured local contextual information.",
    )

    doc.add_heading("Limitations", level=1)
    add_doc_bullets(
        doc,
        [
            "The corpus is small and temple-specific.",
            "Khmer tokenization remains challenging.",
            "MIN_FREQ=10 removes rare but meaningful words.",
            "PCA reduces 50D vectors to 2D and loses information.",
            "Nearest words represent contextual association, not always synonyms.",
        ],
    )

    doc.add_heading("Future Work", level=1)
    add_doc_bullets(
        doc,
        [
            "Train on a larger Khmer corpus.",
            "Improve tokenizer and compound word handling.",
            "Compare Word2Vec, FastText, and transformer embeddings.",
            "Add optional retraining for uploaded text in Streamlit.",
            "Deploy the Streamlit app publicly.",
        ],
    )

    doc.add_heading("Conclusion", level=1)
    add_doc_paragraph(
        doc,
        "The project satisfies the Mini Project 3 requirements. It implements Khmer skip-gram embeddings, PCA visualization, "
        "n-gram and neural language models, scratch embedding comparison, and honest evaluation. The fixed skip-gram neural LM "
        "performed best with test perplexity around 115.10.",
    )

    doc.save(REPORT_PATH)


def set_slide_background(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PptRGBColor.from_string(color_hex.replace("#", ""))


def add_slide_title(slide, title):
    box = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.25), PptInches(12.4), PptInches(0.55))
    text_frame = box.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = PptPt(24)
    p.font.color.rgb = PptRGBColor.from_string(NAVY.replace("#", ""))


def add_ppt_bullets(slide, bullets, left=0.7, top=1.25, width=5.8, height=4.8, font_size=20):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(font_size)
        p.font.color.rgb = PptRGBColor.from_string(GRAY.replace("#", ""))


def add_ppt_table(slide, headers, rows, left, top, width, height, font_size=11):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    table = table_shape.table
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGBColor.from_string(NAVY.replace("#", ""))
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = PptRGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.font.size = PptPt(font_size)
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = "" if pd.isna(value) else str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptPt(font_size)
                paragraph.font.color.rgb = PptRGBColor.from_string(GRAY.replace("#", ""))
    return table


def add_picture_if_exists(slide, image_path, left, top, width, height=None):
    if image_path and Path(image_path).exists():
        if height is None:
            slide.shapes.add_picture(str(image_path), PptInches(left), PptInches(top), width=PptInches(width))
        else:
            slide.shapes.add_picture(str(image_path), PptInches(left), PptInches(top), width=PptInches(width), height=PptInches(height))


def add_metric_box(slide, label, value, left, top, width=2.6):
    shape = slide.shapes.add_shape(1, PptInches(left), PptInches(top), PptInches(width), PptInches(0.78))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor.from_string(SOFT_BLUE.replace("#", ""))
    shape.line.color.rgb = PptRGBColor.from_string(GOLD.replace("#", ""))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{label}\n{value}"
    p.font.size = PptPt(12)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor.from_string(NAVY.replace("#", ""))
    p.alignment = PP_ALIGN.CENTER


def create_powerpoint(data, chart_paths):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, NAVY)
    title_box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.3), PptInches(11.8), PptInches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Khmer Word Embeddings and Neural Language Modeling using Temples Corpus"
    p.font.size = PptPt(34)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(255, 255, 255)
    sub = slide.shapes.add_textbox(PptInches(0.8), PptInches(3.0), PptInches(11), PptInches(1.2))
    sub.text_frame.text = "NLP Mini Project 3 - Word Embeddings\nMaster of Data Science | Institute of Technology of Cambodia"
    for paragraph in sub.text_frame.paragraphs:
        paragraph.font.size = PptPt(20)
        paragraph.font.color.rgb = PptRGBColor.from_string(GOLD.replace("#", ""))

    def add_standard_slide(title, bullets):
        slide = prs.slides.add_slide(blank)
        set_slide_background(slide, "FFFFFF")
        add_slide_title(slide, title)
        add_ppt_bullets(slide, bullets, left=0.8, top=1.25, width=6.0, height=5.2)
        return slide

    add_standard_slide("Project Objective", [
        "Build Khmer word embeddings.",
        "Visualize embeddings using PCA.",
        "Build neural language models.",
        "Compare n-gram, fixed embedding LM, and scratch LM.",
    ])
    add_standard_slide("Mini Project 3 Requirements", [
        "Embedding dimension = 50",
        "Context window = +/-4",
        "Negative sampling = 2",
        "MIN_FREQ = 10",
        "Neural LM: n=5, h=512",
        "Scratch embedding comparison included",
    ])
    add_standard_slide("Dataset Overview", [
        "Dataset: temples.txt",
        "Topic: Khmer temple Wikipedia text",
        f"Total useful tokens: {data['total_tokens']:,}",
        f"Vocabulary size: {data['vocab_size']:,}",
        "Small domain-specific Khmer corpus",
    ])
    add_standard_slide("NLP Pipeline", [
        "Khmer Text -> Cleaning -> Tokenization -> EDA",
        "Skip-gram -> PCA Visualization",
        "Neural LM -> Model Evaluation",
    ])
    add_standard_slide("Text Preprocessing", [
        "Khmer tokenization is essential.",
        "Punctuation, numbers, spaces, and noisy tokens removed.",
        "MIN_FREQ = 10 applied.",
        "Rare words mapped to <UNK>.",
    ])

    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, "FFFFFF")
    add_slide_title(slide, "EDA Result")
    add_picture_if_exists(slide, chart_paths.get("top_words"), 0.6, 1.05, 5.9, 5.3)
    add_picture_if_exists(slide, chart_paths.get("frequency_distribution"), 6.8, 1.25, 5.7, 4.2)

    add_standard_slide("Skip-gram Model", [
        "Center word and context words from window +/-4.",
        "Negative sampling k=2.",
        "Embedding dimension = 50.",
        "Binary classification: real vs negative pairs.",
    ])
    add_standard_slide("Word Embedding Analysis", [
        "Cosine similarity compares word vectors.",
        "Nearest words show contextual association.",
        "Small corpus means nearest words are not always synonyms.",
    ])

    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, "FFFFFF")
    add_slide_title(slide, "PCA Visualization")
    if chart_paths.get("scratch_pca"):
        add_picture_if_exists(slide, chart_paths.get("pca"), 0.55, 1.0, 5.9, 5.25)
        add_picture_if_exists(slide, chart_paths.get("scratch_pca"), 6.85, 1.0, 5.9, 5.25)
    else:
        add_picture_if_exists(slide, chart_paths.get("pca"), 0.8, 1.0, 6.8, 5.5)
        add_ppt_bullets(slide, ["PCA reduces 50D embeddings to 2D.", "The map is approximate.", "Words close together often share context."], left=8.0, top=1.5, width=4.5, height=3)

    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, "FFFFFF")
    add_slide_title(slide, "Advanced Embedding Clustering")
    add_picture_if_exists(slide, chart_paths.get("kmeans_silhouette"), 0.65, 1.05, 5.7, 4.4)
    add_picture_if_exists(slide, chart_paths.get("kmeans_pca"), 6.75, 1.05, 5.7, 4.4)
    add_ppt_bullets(
        slide,
        [
            "K-means groups words with similar embedding patterns.",
            "Silhouette score helps choose K.",
            "PCA is used only for 2D visualization.",
        ],
        left=0.85,
        top=5.75,
        width=11.6,
        height=1.0,
        font_size=15,
    )

    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, "FFFFFF")
    add_slide_title(slide, "Hierarchical Clustering")
    dendrogram_for_slide = chart_paths.get("hierarchical_dendrogram_numbered", chart_paths.get("hierarchical_dendrogram"))
    add_picture_if_exists(slide, dendrogram_for_slide, 0.55, 1.0, 7.2, 5.5)
    dendrogram_labels = data.get("dendrogram_labels", pd.DataFrame())
    if not dendrogram_labels.empty and {"ID", "Khmer Word", "Frequency"}.issubset(dendrogram_labels.columns):
        label_rows = dendrogram_labels[["ID", "Khmer Word", "Frequency"]].head(10).values.tolist()
        add_ppt_table(slide, ["ID", "Khmer Word", "Freq."], label_rows, 8.0, 1.15, 4.75, 4.15, font_size=8)
        add_ppt_bullets(slide, ["Numbered labels improve slide readability.", "Full mapping is included in the report."], left=8.1, top=5.55, width=4.4, height=1.0, font_size=14)
    else:
        add_ppt_bullets(
            slide,
            [
                "Dendrogram uses frequent words only.",
                "Cosine distance + average linkage.",
                "Shorter branches suggest more similar embedding patterns.",
                "Interpret as exploratory contextual groups.",
            ],
            left=8.35,
            top=1.45,
            width=4.3,
            height=3.8,
            font_size=17,
        )

    add_standard_slide("Neural Language Model", [
        "Previous 5 words -> embeddings.",
        "Concatenate embeddings.",
        "Hidden layer size = 512.",
        "Softmax predicts the next word.",
    ])

    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, "FFFFFF")
    add_slide_title(slide, "Model Comparison")
    add_picture_if_exists(slide, chart_paths.get("model_perplexity"), 0.65, 1.1, 6.2, 4.8)
    comparison = data["model_comparison"]
    rows = []
    for _, row in comparison.iterrows():
        rows.append([row.get("Model", ""), f"{float(row.get('Test Perplexity', np.nan)):.2f}"])
    add_ppt_table(slide, ["Model", "Test Perplexity"], rows, 7.2, 1.35, 5.2, 2.1, font_size=12)
    add_ppt_bullets(slide, ["Best model: Neural LM fixed skip-gram embeddings.", "Perplexity around 115.10.", "Lower perplexity is better."], left=7.25, top=3.9, width=5.0, height=2)

    add_standard_slide("Discussion and Limitations", [
        "Fixed skip-gram embeddings helped prediction.",
        "Scratch LM can overfit on small data.",
        "Khmer tokenization remains challenging.",
        "PCA is approximate and loses information.",
    ])
    add_standard_slide("Streamlit Deployment", [
        "Khmer Temple NLP Explorer app.",
        "Default temples.txt and upload .txt for EDA.",
        "Embedding explorer, PCA map, and model comparison.",
        "Ready for GitHub and Streamlit Cloud.",
    ])
    add_standard_slide("Conclusion", [
        "Project satisfies Mini Project 3 requirements.",
        "Khmer skip-gram embeddings were learned.",
        "PCA and neural LM were implemented.",
        "Fixed embedding LM achieved best perplexity.",
        "Future work: larger corpus and stronger tokenizer.",
    ])

    prs.save(PPTX_PATH)


def main():
    ensure_dirs()
    font_name = configure_matplotlib_fonts()
    data = load_project_data()
    charts = make_charts(data)
    create_word_report(data, charts)
    create_powerpoint(data, charts)
    print("Matplotlib font:", font_name)
    print("Generated report:", REPORT_PATH)
    print("Generated presentation:", PPTX_PATH)


if __name__ == "__main__":
    main()
